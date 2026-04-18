"""
EAA Document Tools - Phase 2
PDF, Word (docx), Excel (xlsx), PowerPoint (pptx) reading and creation.
All tools use the existing ToolResult pattern from eaa_agent_tools.py.
"""

import os
import json
import io
import traceback
from typing import Optional, List, Dict, Any
from dataclasses import dataclass
from datetime import datetime

try:
    from eaa_agent_tools import ToolResult
except ImportError:
    @dataclass
    class ToolResult:
        success: bool
        output: str
        error: Optional[str] = None
        def to_dict(self):
            return {"success": self.success, "output": self.output, "error": self.error}


# ═══════════════════════════════════════════════════════════════════════════════
# PDF TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def tool_pdf_read(file_path: str, page: int = None, max_chars: int = 10000) -> ToolResult:
    """Read text from a PDF file. Use page param for specific page, or omit for all pages."""
    try:
        file_path = os.path.expanduser(file_path)
        if not os.path.exists(file_path):
            return ToolResult(False, "", f"PDF not found: {file_path}")

        import fitz  # PyMuPDF

        doc = fitz.open(file_path)

        if page is not None:
            if page < 1 or page > len(doc):
                return ToolResult(False, "", f"Page {page} out of range (1-{len(doc)})")
            text = doc[page - 1].get_text()
            return ToolResult(True, f"--- Page {page}/{len(doc)} ---\n{text}")
        else:
            full_text = []
            for i, pg in enumerate(doc):
                text = pg.get_text().strip()
                if text:
                    full_text.append(f"--- Page {i+1}/{len(doc)} ---\n{text}")
            result = "\n\n".join(full_text)
            if len(result) > max_chars:
                result = result[:max_chars] + f"\n\n... [truncated, {len(result)} total chars]"
            return ToolResult(True, f"PDF: {file_path} ({len(doc)} pages)\n\n{result}")

    except ImportError:
        # Fallback to pypdf
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            pages = []
            for i, pg in enumerate(reader.pages):
                text = pg.extract_text()
                if text:
                    pages.append(f"--- Page {i+1} ---\n{text}")
            result = "\n\n".join(pages)
            return ToolResult(True, f"PDF: {file_path} ({len(reader.pages)} pages)\n\n{result[:max_chars]}")
        except Exception as e2:
            return ToolResult(False, "", f"PDF read failed: {e2}. Install: pip install PyMuPDF")
    except Exception as e:
        return ToolResult(False, "", f"PDF read failed: {str(e)}")


def tool_pdf_info(file_path: str) -> ToolResult:
    """Get PDF metadata and structure information."""
    try:
        file_path = os.path.expanduser(file_path)
        import fitz
        doc = fitz.open(file_path)

        metadata = doc.metadata or {}
        info = {
            "file": file_path,
            "pages": len(doc),
            "file_size": f"{os.path.getsize(file_path):,} bytes",
            "title": metadata.get("title", "N/A"),
            "author": metadata.get("author", "N/A"),
            "subject": metadata.get("subject", "N/A"),
            "creator": metadata.get("creator", "N/A"),
            "creation_date": metadata.get("creationDate", "N/A"),
            "modification_date": metadata.get("modDate", "N/A"),
            "pdf_version": metadata.get("format", "N/A"),
            "encrypted": doc.is_encrypted,
        }

        # Page dimensions
        if len(doc) > 0:
            first_page = doc[0]
            info["first_page_size"] = f"{first_page.rect.width:.0f}x{first_page.rect.height:.0f} points"

        doc.close()
        return ToolResult(True, json.dumps(info, indent=2))

    except Exception as e:
        return ToolResult(False, "", f"PDF info failed: {str(e)}")


def tool_pdf_create(
    file_path: str,
    title: str = "",
    content: str = "",
    pages: str = None
) -> ToolResult:
    """
    Create a PDF document.
    content: Single block of text for the whole document.
    pages: JSON string of list, e.g. '[{"title": "Page 1", "text": "Hello"}, ...]'
    """
    try:
        file_path = os.path.expanduser(file_path)
        os.makedirs(os.path.dirname(file_path) if os.path.dirname(file_path) else ".", exist_ok=True)

        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import inch

        doc = SimpleDocTemplate(file_path, pagesize=letter,
                                title=title or "EAA Generated PDF",
                                author="EAA - Embedded AI Assistant")
        styles = getSampleStyleSheet()
        story = []

        # Title
        if title:
            story.append(Paragraph(title, styles["Title"]))
            story.append(Spacer(1, 0.3 * inch))

        if pages:
            # Multi-page mode
            page_list = json.loads(pages) if isinstance(pages, str) else pages
            for pg in page_list:
                pg_title = pg.get("title", "")
                pg_text = pg.get("text", "")
                if pg_title:
                    story.append(Paragraph(pg_title, styles["Heading2"]))
                if pg_text:
                    # Split into paragraphs by double newline
                    for para in pg_text.split("\n\n"):
                        if para.strip():
                            story.append(Paragraph(para.strip().replace("\n", "<br/>"), styles["Normal"]))
                story.append(Spacer(1, 0.2 * inch))
        elif content:
            # Single content mode
            for para in content.split("\n\n"):
                if para.strip():
                    story.append(Paragraph(para.strip().replace("\n", "<br/>"), styles["Normal"]))

        doc.build(story)
        size = os.path.getsize(file_path)
        return ToolResult(True, f"PDF created: {file_path} ({size:,} bytes)")

    except Exception as e:
        return ToolResult(False, "", f"PDF creation failed: {str(e)}")


# ═══════════════════════════════════════════════════════════════════════════════
# DOCX (WORD) TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def tool_docx_read(file_path: str, max_chars: int = 10000) -> ToolResult:
    """Read text content from a Word (.docx) file."""
    try:
        file_path = os.path.expanduser(file_path)
        if not os.path.exists(file_path):
            return ToolResult(False, "", f"DOCX not found: {file_path}")

        from docx import Document
        doc = Document(file_path)

        # Extract paragraphs
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Extract tables
        tables_text = []
        for i, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text for cell in row.cells]
                rows.append(" | ".join(cells))
            tables_text.append(f"--- Table {i+1} ---\n" + "\n".join(rows))

        result = "\n\n".join(paragraphs)
        if tables_text:
            result += "\n\n" + "\n\n".join(tables_text)

        if len(result) > max_chars:
            result = result[:max_chars] + f"\n\n... [truncated, {len(result)} total chars]"

        return ToolResult(True, f"DOCX: {file_path} ({len(paragraphs)} paragraphs, {len(doc.tables)} tables)\n\n{result}")

    except Exception as e:
        return ToolResult(False, "", f"DOCX read failed: {str(e)}")


def tool_docx_create(
    file_path: str,
    title: str = "",
    content: str = "",
    paragraphs: str = None
) -> ToolResult:
    """
    Create a Word (.docx) document.
    content: Full text content (split by \\n\\n for paragraphs).
    paragraphs: JSON string of list of dicts, e.g. '[{"text": "Hello", "style": "Heading 1"}, ...]'
    styles: Normal, Heading 1-3, Title, Subtitle, List Bullet, Quote
    """
    try:
        file_path = os.path.expanduser(file_path)
        os.makedirs(os.path.dirname(file_path) if os.path.dirname(file_path) else ".", exist_ok=True)

        from docx import Document
        from docx.shared import Pt, Inches

        doc = Document()

        if title:
            doc.add_heading(title, level=0)

        if paragraphs:
            para_list = json.loads(paragraphs) if isinstance(paragraphs, str) else paragraphs
            for p in para_list:
                text = p.get("text", "")
                style = p.get("style", "Normal")
                if style.startswith("Heading"):
                    level = int(style.split()[-1]) if style.split()[-1].isdigit() else 1
                    doc.add_heading(text, level=level)
                elif style == "Title":
                    doc.add_heading(text, level=0)
                elif style == "Subtitle":
                    doc.add_heading(text, level=2)
                else:
                    doc.add_paragraph(text, style=style)
        elif content:
            for para in content.split("\n\n"):
                if para.strip():
                    doc.add_paragraph(para.strip())

        doc.save(file_path)
        size = os.path.getsize(file_path)
        return ToolResult(True, f"DOCX created: {file_path} ({size:,} bytes)")

    except Exception as e:
        return ToolResult(False, "", f"DOCX creation failed: {str(e)}")


# ═══════════════════════════════════════════════════════════════════════════════
# XLSX (EXCEL) TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def tool_xlsx_read(file_path: str, sheet_name: str = None, max_rows: int = 100) -> ToolResult:
    """Read data from an Excel (.xlsx) file."""
    try:
        file_path = os.path.expanduser(file_path)
        if not os.path.exists(file_path):
            return ToolResult(False, "", f"XLSX not found: {file_path}")

        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)

        # Sheet info
        sheets = wb.sheetnames
        output_parts = [f"XLSX: {file_path} ({len(sheets)} sheets: {', '.join(sheets)})"]

        target_sheet = sheet_name or sheets[0]
        if target_sheet not in sheets:
            return ToolResult(False, "", f"Sheet '{target_sheet}' not found. Available: {', '.join(sheets)}")

        ws = wb[target_sheet]

        rows_data = []
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            if row_count >= max_rows:
                rows_data.append(f"... ({max_rows} rows shown, more available)")
                break
            cells = [str(c) if c is not None else "" for c in row]
            rows_data.append(" | ".join(cells))
            row_count += 1

        output_parts.append(f"\n--- Sheet: {target_sheet} ({ws.max_row} rows x {ws.max_column} cols) ---")
        output_parts.append("\n".join(rows_data))

        wb.close()
        return ToolResult(True, "\n".join(output_parts))

    except Exception as e:
        return ToolResult(False, "", f"XLSX read failed: {str(e)}")


def tool_xlsx_create(
    file_path: str,
    sheet_name: str = "Sheet1",
    headers: str = None,
    rows: str = None
) -> ToolResult:
    """
    Create an Excel (.xlsx) file.
    headers: JSON array of column names, e.g. '["Name", "Age", "City"]'
    rows: JSON array of arrays, e.g. '[["Alice", 30, "NYC"], ["Bob", 25, "LA"]]'
    """
    try:
        file_path = os.path.expanduser(file_path)
        os.makedirs(os.path.dirname(file_path) if os.path.dirname(file_path) else ".", exist_ok=True)

        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill

        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name

        if headers:
            header_list = json.loads(headers) if isinstance(headers, str) else headers
            # Style headers
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF")
            for col, header in enumerate(header_list, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center")

        if rows:
            row_list = json.loads(rows) if isinstance(rows, str) else rows
            start_row = 2 if headers else 1
            for r_idx, row_data in enumerate(row_list, start_row):
                for c_idx, value in enumerate(row_data, 1):
                    ws.cell(row=r_idx, column=c_idx, value=value)

        # Auto-fit column widths
        for column in ws.columns:
            max_length = 0
            for cell in column:
                try:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                except Exception:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column[0].column_letter].width = adjusted_width

        wb.save(file_path)
        size = os.path.getsize(file_path)
        return ToolResult(True, f"XLSX created: {file_path} ({size:,} bytes, sheet: {sheet_name})")

    except Exception as e:
        return ToolResult(False, "", f"XLSX creation failed: {str(e)}")


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX (POWERPOINT) TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def tool_pptx_read(file_path: str, max_chars: int = 10000) -> ToolResult:
    """Read text content from a PowerPoint (.pptx) file."""
    try:
        file_path = os.path.expanduser(file_path)
        if not os.path.exists(file_path):
            return ToolResult(False, "", f"PPTX not found: {file_path}")

        from pptx import Presentation
        prs = Presentation(file_path)

        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            # Check if it's a title
                            level = para.level
                            prefix = "  " * level + ("- " if level > 0 else "")
                            slide_content.append(f"{prefix}{para.text.strip()}")
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text for cell in row.cells]
                        slide_content.append(" | ".join(cells))

            if slide_content:
                slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_content))

        result = f"PPTX: {file_path} ({len(prs.slides)} slides)\n\n" + "\n\n".join(slides_text)

        if len(result) > max_chars:
            result = result[:max_chars] + f"\n\n... [truncated]"

        return ToolResult(True, result)

    except Exception as e:
        return ToolResult(False, "", f"PPTX read failed: {str(e)}")


def tool_pptx_create(
    file_path: str,
    title: str = "",
    slides: str = None,
    content: str = ""
) -> ToolResult:
    """
    Create a PowerPoint (.pptx) file.
    slides: JSON array of dicts, e.g. '[{"title": "Slide 1", "content": "Hello world", "layout": "title_and_content"}, ...]'
    content: Simple text, each slide separated by '--- SLIDE ---'
    layouts: title_slide, title_and_content, blank, two_content
    """
    try:
        file_path = os.path.expanduser(file_path)
        os.makedirs(os.path.dirname(file_path) if os.path.dirname(file_path) else ".", exist_ok=True)

        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()

        if slides:
            slide_list = json.loads(slides) if isinstance(slides, str) else slides
            for slide_data in slide_list:
                slide_title = slide_data.get("title", "")
                slide_content = slide_data.get("content", "")
                layout_name = slide_data.get("layout", "title_and_content")

                # Choose layout
                if layout_name == "title_slide":
                    layout = prs.slide_layouts[0]  # Title Slide
                    slide = prs.slides.add_slide(layout)
                    if slide_title:
                        slide.shapes.title.text = slide_title
                    if slide_content:
                        if slide.placeholders[1]:
                            slide.placeholders[1].text = slide_content
                elif layout_name == "blank":
                    layout = prs.slide_layouts[6]  # Blank
                    slide = prs.slides.add_slide(layout)
                    if slide_title:
                        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                        tf = txBox.text_frame
                        p = tf.paragraphs[0]
                        p.text = slide_title
                        p.font.size = Pt(32)
                        p.font.bold = True
                    if slide_content:
                        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        for line in slide_content.split("\n"):
                            p = tf.add_paragraph()
                            p.text = line
                            p.font.size = Pt(18)
                else:
                    layout = prs.slide_layouts[1]  # Title and Content
                    slide = prs.slides.add_slide(layout)
                    if slide_title:
                        slide.shapes.title.text = slide_title
                    if slide_content:
                        tf = slide.placeholders[1].text_frame
                        tf.clear()
                        for line in slide_content.split("\n"):
                            p = tf.add_paragraph()
                            p.text = line
                            p.font.size = Pt(18)
        elif content:
            # Simple mode: split by slide separator
            slide_parts = content.split("--- SLIDE ---")
            for part in slide_parts:
                part = part.strip()
                if not part:
                    continue
                lines = part.split("\n")
                slide_title = lines[0].strip() if lines else ""
                slide_body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""

                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_title
                if slide_body:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    for line in slide_body.split("\n"):
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = Pt(18)

        prs.save(file_path)
        size = os.path.getsize(file_path)
        return ToolResult(True, f"PPTX created: {file_path} ({size:,} bytes, {len(prs.slides)} slides)")

    except Exception as e:
        return ToolResult(False, "", f"PPTX creation failed: {str(e)}")


# ═══════════════════════════════════════════════════════════════════════════════
# REGISTRY
# ═══════════════════════════════════════════════════════════════════════════════

def register_document_tools(registry) -> None:
    """Register all document tools with the existing ToolRegistry."""
    registry.register("pdf_read", tool_pdf_read, "Read PDF text. Args: file_path, page (optional), max_chars")
    registry.register("pdf_info", tool_pdf_info, "Get PDF metadata. Args: file_path")
    registry.register("pdf_create", tool_pdf_create, "Create PDF document. Args: file_path, title, content or pages (JSON)")
    registry.register("docx_read", tool_docx_read, "Read Word document. Args: file_path, max_chars")
    registry.register("docx_create", tool_docx_create, "Create Word document. Args: file_path, title, content or paragraphs (JSON)")
    registry.register("xlsx_read", tool_xlsx_read, "Read Excel spreadsheet. Args: file_path, sheet_name (optional), max_rows")
    registry.register("xlsx_create", tool_xlsx_create, "Create Excel spreadsheet. Args: file_path, headers (JSON array), rows (JSON array)")
    registry.register("pptx_read", tool_pptx_read, "Read PowerPoint. Args: file_path, max_chars")
    registry.register("pptx_create", tool_pptx_create, "Create PowerPoint. Args: file_path, title, slides (JSON) or content")

__all__ = [
    "register_document_tools",
    "tool_pdf_read", "tool_pdf_info", "tool_pdf_create",
    "tool_docx_read", "tool_docx_create",
    "tool_xlsx_read", "tool_xlsx_create",
    "tool_pptx_read", "tool_pptx_create",
]
